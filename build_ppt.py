"""
SalesCast College Presentation
Theme: Dark navy/teal matching the website (#0a0d14, #00e5b0, #3d8bff)
9 slides, college format
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours (website theme) ──────────────────────────────
BG_DARK    = RGBColor(0x0A, 0x0D, 0x14)   # #0a0d14  main bg
BG_CARD    = RGBColor(0x11, 0x16, 0x22)   # #111622  card
BG_CARD2   = RGBColor(0x1A, 0x21, 0x33)   # #1a2133  lighter card
ACCENT     = RGBColor(0x00, 0xE5, 0xB0)   # #00e5b0  teal
ACCENT2    = RGBColor(0x3D, 0x8B, 0xFF)   # #3d8bff  blue
WARN       = RGBColor(0xFF, 0x6B, 0x35)   # #ff6b35  orange
TEXT_BRT   = RGBColor(0xF0, 0xF6, 0xFF)   # bright white
TEXT_DIM   = RGBColor(0x5A, 0x6D, 0x8A)   # muted
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BORDER     = RGBColor(0x1E, 0x2D, 0x45)

W  = Inches(13.33)   # LAYOUT_WIDE
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(bg_color=BG_DARK):
    layout = prs.slide_layouts[6]   # blank
    sl = prs.slides.add_slide(layout)
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return sl

def rect(sl, x, y, w, h, color, alpha=None):
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def txt(sl, text, x, y, w, h, size=16, bold=False, color=TEXT_BRT,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.color.rgb = color
    run.font.italic  = italic
    run.font.name    = font
    return tb

def txts(sl, lines, x, y, w, h, size=14, color=TEXT_BRT,
         align=PP_ALIGN.LEFT, bold=False, font="Calibri", spacing=1.15):
    """Multi-line textbox from list of strings"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = font
    return tb

def accent_bar(sl, x, y, h, color=ACCENT, w=0.06):
    """Thin vertical accent bar"""
    rect(sl, x, y, w, h, color)

def slide_num(sl, n, total=9):
    txt(sl, f"{n} / {total}", 12.4, 7.1, 0.8, 0.3,
        size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def footer_bar(sl, label="SalesCast — ML Mini Project"):
    rect(sl, 0, 7.18, 13.33, 0.32, BG_CARD)
    txt(sl, label, 0.4, 7.2, 8, 0.28, size=9, color=TEXT_DIM)
    txt(sl, "Linear Regression | Python · Scikit-learn · Flask",
        7, 7.2, 6.2, 0.28, size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def section_tag(sl, tag_text, x=0.4, y=0.28):
    rect(sl, x, y, len(tag_text)*0.095 + 0.25, 0.3, BG_CARD2)
    txt(sl, tag_text, x+0.08, y+0.03, 2, 0.25,
        size=9, color=ACCENT, bold=True, font="Consolas")

def card(sl, x, y, w, h):
    shp = sl.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_CARD
    shp.line.color.rgb = BORDER
    shp.line.width = Emu(12700)   # 1pt
    return shp

def stat_block(sl, x, y, num, label, color=ACCENT):
    card(sl, x, y, 2.6, 1.3)
    txt(sl, num, x+0.15, y+0.08, 2.3, 0.65,
        size=32, bold=True, color=color, font="Calibri")
    txt(sl, label, x+0.15, y+0.72, 2.3, 0.45,
        size=11, color=TEXT_DIM, font="Calibri")

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / COVER
# ═══════════════════════════════════════════════════════════
s1 = blank_slide()

# Space for college banner (top placeholder)
rect(s1, 0, 0, 13.33, 1.7, BG_CARD)
shp = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.7))
shp.fill.solid(); shp.fill.fore_color.rgb = BG_CARD2
shp.line.fill.background()
txt(s1, "[ College Banner / Logo — paste your image here ]",
    0, 0.6, 13.33, 0.5, size=13, color=TEXT_DIM,
    align=PP_ALIGN.CENTER, italic=True)

# Teal left accent strip
rect(s1, 0, 1.7, 0.5, 5.48, ACCENT)

# Title block
txt(s1, "SalesCast", 0.75, 1.95, 12, 1.05,
    size=56, bold=True, color=TEXT_BRT, font="Calibri")
txt(s1, "Sales Forecasting Using Linear Regression",
    0.75, 2.95, 11, 0.6, size=24, color=ACCENT, font="Calibri")
txt(s1, "A Machine Learning Mini Project",
    0.75, 3.5, 10, 0.45, size=16, color=TEXT_DIM, italic=True)

# Divider
rect(s1, 0.75, 4.05, 10.0, 0.04, BORDER)

# Team + Guide — side by side
rect(s1, 0.75, 4.25, 5.8, 1.7, BG_CARD)
rect(s1, 0.75, 4.25, 0.06, 1.7, ACCENT2)
txt(s1, "Team Members", 0.97, 4.32, 5.4, 0.35,
    size=11, bold=True, color=ACCENT2)
txts(s1, ["Student Name 1  —  Roll No.",
          "Student Name 2  —  Roll No.",
          "Student Name 3  —  Roll No."],
     0.97, 4.65, 5.4, 1.15, size=13, color=TEXT_BRT)

rect(s1, 7.0, 4.25, 5.8, 1.7, BG_CARD)
rect(s1, 7.0, 4.25, 0.06, 1.7, ACCENT)
txt(s1, "Project Guide", 7.22, 4.32, 5.4, 0.35,
    size=11, bold=True, color=ACCENT)
txts(s1, ["Prof. / Dr. Guide Name",
          "Department of Computer Science",
          "Academic Year 2025–26"],
     7.22, 4.65, 5.4, 1.15, size=13, color=TEXT_BRT)

footer_bar(s1)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT & OBJECTIVE
# ═══════════════════════════════════════════════════════════
s2 = blank_slide()
section_tag(s2, "02  PROBLEM STATEMENT")
txt(s2, "What Problem Are We Solving?",
    0.4, 0.62, 12, 0.7, size=34, bold=True, color=TEXT_BRT)
footer_bar(s2); slide_num(s2, 2)

# Left: problem
card(s2, 0.4, 1.5, 5.8, 4.8)
accent_bar(s2, 0.4, 1.5, 4.8, WARN)
txt(s2, "The Problem", 0.65, 1.62, 5.3, 0.4, size=14, bold=True, color=WARN)
txts(s2, [
    "Businesses make ad spend and hiring decisions blind.",
    "",
    "Without data-driven forecasts, they either overspend",
    "on marketing that won't return, or understaff during",
    "peak demand periods.",
    "",
    "Sales managers rely on gut instinct instead of",
    "patterns that already exist in their own data."
], 0.65, 2.1, 5.2, 3.9, size=13.5, color=TEXT_BRT)

# Right: objective
card(s2, 6.7, 1.5, 6.2, 4.8)
accent_bar(s2, 6.7, 1.5, 4.8, ACCENT)
txt(s2, "Our Objective", 6.95, 1.62, 5.7, 0.4, size=14, bold=True, color=ACCENT)
txts(s2, [
    "Build a web-based ML system that:",
    "",
    "→  Takes 3 business inputs as features",
    "     (Ad Spend, Salespeople, Price per Unit)",
    "",
    "→  Uses Linear Regression to learn the",
    "     relationship between inputs and Sales",
    "",
    "→  Predicts monthly sales revenue in",
    "     both USD ($) and INR (₹)"
], 6.95, 2.1, 5.7, 3.9, size=13.5, color=TEXT_BRT)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — LINEAR REGRESSION THEORY
# ═══════════════════════════════════════════════════════════
s3 = blank_slide()
section_tag(s3, "03  ALGORITHM")
txt(s3, "Linear Regression — How It Works",
    0.4, 0.62, 12, 0.7, size=34, bold=True, color=TEXT_BRT)
footer_bar(s3); slide_num(s3, 3)

# Equation box
rect(s3, 0.4, 1.5, 12.53, 0.95, BG_CARD)
rect(s3, 0.4, 1.5, 12.53, 0.04, ACCENT)
txt(s3, "Sales  =  β₀  +  β₁(Ad Spend)  +  β₂(Salespeople)  +  β₃(Avg Price)",
    0.6, 1.58, 12.1, 0.75, size=19, bold=True, color=ACCENT, font="Consolas")

# Three concept cards
for i, (title, body, clr) in enumerate([
    ("β₀  — Intercept",
     "The baseline sales value when all features are zero. Our model gives β₀ = 63,812. This is the model's 'starting point' before any input is considered.",
     ACCENT2),
    ("Positive Coefficient",
     "β₁ (Ad Spend) = +21,728  and  β₂ (Salespeople) = +9,688. Positive means — as these inputs increase, predicted Sales goes UP. More spend, more people → higher revenue.",
     ACCENT),
    ("Negative Coefficient",
     "β₃ (Avg Price) = −5,442. Negative means — as price rises, fewer units are sold, so total revenue DROPS. This is the classic price-demand relationship from economics.",
     WARN),
]):
    cx = 0.4 + i * 4.3
    card(s3, cx, 2.65, 4.1, 3.65)
    accent_bar(s3, cx, 2.65, 3.65, clr)
    txt(s3, title, cx+0.25, 2.78, 3.7, 0.45, size=13, bold=True, color=clr)
    txts(s3, body.split(". "), cx+0.25, 3.28, 3.7, 2.9, size=12, color=TEXT_BRT)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — DATASET
# ═══════════════════════════════════════════════════════════
s4 = blank_slide()
section_tag(s4, "04  DATASET")
txt(s4, "Training Data — What We Fed the Model",
    0.4, 0.62, 12, 0.7, size=34, bold=True, color=TEXT_BRT)
footer_bar(s4); slide_num(s4, 4)

# Stats row
stat_block(s4, 0.4,  1.55, "50",      "Data samples (rows)",    ACCENT)
stat_block(s4, 3.2,  1.55, "3",       "Input features (X)",     ACCENT2)
stat_block(s4, 6.0,  1.55, "1",       "Target variable (Y)",    ACCENT)
stat_block(s4, 8.8,  1.55, "80 / 20", "Train / Test split",     ACCENT2)
stat_block(s4, 11.6 - 2.7, 1.55, "R² ≈ 0.99", "Model accuracy", WARN)

# Feature table
headers = ["Feature (X)", "Example Values", "Effect on Sales"]
rows = [
    ["Advertising Spend (₹)", "₹2,50,500 → ₹11,69,000",  "Higher spend → more customers → Sales UP  ▲"],
    ["Num. Salespeople",      "8 → 21 people",     "More reps → more deals closed → Sales UP  ▲"],
    ["Avg Price per Unit (₹)","₹1,670 → ₹3,757",          "Higher price → less demand → Sales DOWN  ▼"],
    ["Sales Revenue (₹)",     "₹23,38,000 → ₹1,06,88,000","Target — what the model learns to predict"],
]

# Table header row
hy = 3.1
rect(s4, 0.4, hy, 12.53, 0.42, BG_CARD2)
for ci, (hdr, cw, cx) in enumerate(zip(headers, [3.5, 3.5, 5.5],
                                        [0.4, 3.9, 7.4])):
    txt(s4, hdr, cx+0.1, hy+0.08, cw-0.1, 0.3,
        size=12, bold=True, color=ACCENT, font="Consolas")

for ri, row in enumerate(rows):
    ry = hy + 0.42 + ri * 0.82
    bg = BG_CARD if ri % 2 == 0 else BG_CARD2
    rect(s4, 0.4, ry, 12.53, 0.8, bg)
    clr = WARN if "DOWN" in row[2] else TEXT_BRT
    for ci, (cell, cw, cx) in enumerate(zip(row, [3.5, 3.5, 5.5],
                                              [0.4, 3.9, 7.4])):
        c = clr if ci == 2 else TEXT_BRT
        txt(s4, cell, cx+0.1, ry+0.18, cw-0.1, 0.5, size=12, color=c)

rect(s4, 0.4, hy, 12.53, 0.04, ACCENT)   # top accent line on table

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — METHODOLOGY / SYSTEM FLOW
# ═══════════════════════════════════════════════════════════
s5 = blank_slide()
section_tag(s5, "05  METHODOLOGY")
txt(s5, "How the System Works — Step by Step",
    0.4, 0.62, 12, 0.7, size=34, bold=True, color=TEXT_BRT)
footer_bar(s5); slide_num(s5, 5)

steps = [
    ("01", "Collect Data",       "50 rows of real sales records with 3 features", ACCENT),
    ("02", "Pre-process",        "StandardScaler normalises features so no single variable dominates", ACCENT2),
    ("03", "Split Dataset",      "80% used for training, 20% held back for testing accuracy", ACCENT),
    ("04", "Train Model",        "scikit-learn LinearRegression fits the best line through training data", ACCENT2),
    ("05", "Evaluate",           "R², RMSE, MAE computed on test set — R² = 0.99 achieved", ACCENT),
    ("06", "Predict & Deploy",   "Flask API serves predictions through a browser UI in real-time", WARN),
]

# Two columns of 3
for i, (num, title, body, clr) in enumerate(steps):
    col = i % 2
    row = i // 2
    cx = 0.4 + col * 6.5
    cy = 1.55 + row * 1.85
    card(s5, cx, cy, 6.1, 1.65)
    # Number circle
    rect(s5, cx+0.15, cy+0.35, 0.55, 0.55, clr)
    txt(s5, num, cx+0.15, cy+0.38, 0.55, 0.48,
        size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER, font="Consolas")
    txt(s5, title, cx+0.85, cy+0.2, 5.0, 0.42,
        size=14, bold=True, color=clr)
    txt(s5, body, cx+0.85, cy+0.62, 5.0, 0.85,
        size=12, color=TEXT_BRT)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — TECH STACK
# ═══════════════════════════════════════════════════════════
s6 = blank_slide()
section_tag(s6, "06  TECH STACK")
txt(s6, "Tools & Technologies Used",
    0.4, 0.62, 12, 0.7, size=34, bold=True, color=TEXT_BRT)
footer_bar(s6); slide_num(s6, 6)

techs = [
    ("Python 3",      "Core programming language. Simple syntax, vast ML library support.", ACCENT),
    ("scikit-learn",  "LinearRegression, StandardScaler, train_test_split, evaluation metrics.", ACCENT2),
    ("Flask",         "Lightweight web framework — serves the API and renders the HTML UI.", ACCENT),
    ("pandas / NumPy","Data loading, cleaning, manipulation and numerical operations.", ACCENT2),
    ("joblib",        "Saves and loads the trained model (.pkl) so retraining is not needed on every launch.", ACCENT),
    ("Chart.js",      "Browser-side library that draws the Coefficients and Trend charts on the web UI.", ACCENT2),
    ("HTML / CSS / JS","Complete responsive dark-themed dashboard for prediction and model insight.", WARN),
    ("VS Code",       "Development environment — all code written, tested and debugged here.", TEXT_DIM),
]

for i, (name, desc, clr) in enumerate(techs):
    col = i % 2
    row = i // 2
    cx = 0.4 + col * 6.5
    cy = 1.5 + row * 1.38
    card(s6, cx, cy, 6.1, 1.2)
    accent_bar(s6, cx, cy, 1.2, clr)
    txt(s6, name, cx+0.25, cy+0.12, 5.7, 0.38,
        size=14, bold=True, color=clr)
    txt(s6, desc, cx+0.25, cy+0.52, 5.7, 0.6,
        size=11.5, color=TEXT_BRT)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — RESULTS & MODEL PERFORMANCE
# ═══════════════════════════════════════════════════════════
s7 = blank_slide()
section_tag(s7, "07  RESULTS")
txt(s7, "Model Performance & Key Results",
    0.4, 0.62, 12, 0.7, size=34, bold=True, color=TEXT_BRT)
footer_bar(s7); slide_num(s7, 7)

# Big metrics
stat_block(s7, 0.4,  1.55, "0.99",    "R² Score (99% accuracy)",   ACCENT)
stat_block(s7, 3.15, 1.55, "₹2,37,975",  "RMSE (Avg. Error)",          ACCENT2)
stat_block(s7, 5.9,  1.55, "₹1,75,350",  "MAE (Mean Abs. Error)",      ACCENT)
stat_block(s7, 8.65, 1.55, "50",      "Training samples",           ACCENT2)
stat_block(s7, 11.4 - 2.55, 1.55, "10", "Test samples (20%)",       WARN)

# Sample predictions table
rect(s7, 0.4, 3.05, 12.53, 0.42, BG_CARD2)
rect(s7, 0.4, 3.05, 12.53, 0.04, ACCENT)
for ci, (h, cw, cx) in enumerate(zip(
    ["Ad Spend (₹)", "Salespeople", "Avg Price (₹)", "Actual Sales (₹)", "Predicted Sales (₹)", "Error (₹)"],
    [2.0, 1.8, 2.0, 2.3, 2.5, 1.93],
    [0.4, 2.4, 4.2, 6.2, 8.5, 11.0]
)):
    txt(s7, h, cx+0.08, 3.12, cw, 0.28,
        size=11, bold=True, color=ACCENT, font="Consolas")

preds = [
    ["4,17,500",  "10", "2,087", "37,57,500",  "37,40,800",  "16,700"],
    ["8,35,000", "15", "2,922", "74,31,500",  "73,98,100",  "33,400"],
    ["11,27,250", "20", "3,590", "1,00,20,000","99,69,900","50,100"],
    ["6,26,250",  "12", "2,505", "51,77,000",  "51,56,125",  "20,875"],
]
for ri, row in enumerate(preds):
    ry = 3.47 + ri * 0.82
    rect(s7, 0.4, ry, 12.53, 0.78, BG_CARD if ri%2==0 else BG_CARD2)
    for ci, (cell, cw, cx) in enumerate(zip(row,
        [2.0, 1.8, 2.0, 2.3, 2.5, 1.93],
        [0.4, 2.4, 4.2, 6.2, 8.5, 11.0]
    )):
        clr = ACCENT if ci == 4 else TEXT_BRT
        txt(s7, cell, cx+0.08, ry+0.2, cw, 0.42, size=12, color=clr)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — WEB UI SCREENSHOTS (described as wireframes)
# ═══════════════════════════════════════════════════════════
s8 = blank_slide()
section_tag(s8, "08  WEB APPLICATION")
txt(s8, "SalesCast — Web Dashboard",
    0.4, 0.62, 12, 0.7, size=34, bold=True, color=TEXT_BRT)
footer_bar(s8); slide_num(s8, 8)

features = [
    ("⚡ Predict Sales",
     "Enter ad spend, number of salespeople, and average unit price. Get instant sales prediction in both USD and INR with a single click.",
     ACCENT),
    ("₹ / $ Currency Toggle",
     "A toggle switch lets users switch between USD and INR inputs. Sliders auto-recalibrate. The result always shows both currencies.",
     ACCENT2),
    ("📈 Model Visualization",
     "Coefficients bar chart shows impact of each feature. Trend scatter plot shows the linear relationship learned by the model.",
     ACCENT),
    ("🔬 Model Equation",
     "The full equation is shown with plain-English breakdown — every term explained so non-technical users can understand the prediction.",
     ACCENT2),
    ("🗂️ Dataset Preview",
     "All 50 training rows displayed in a scrollable table so the audience can verify the data the model was trained on.",
     WARN),
    ("🔄 Retrain Button",
     "One-click retraining reruns the entire ML pipeline on the current dataset and updates the model file and all displayed metrics.",
     ACCENT),
]

for i, (title, desc, clr) in enumerate(features):
    col = i % 2
    row = i // 2
    cx = 0.4 + col * 6.5
    cy = 1.5 + row * 1.72
    card(s8, cx, cy, 6.1, 1.55)
    accent_bar(s8, cx, cy, 1.55, clr)
    txt(s8, title, cx+0.25, cy+0.12, 5.7, 0.42, size=13, bold=True, color=clr)
    txt(s8, desc,  cx+0.25, cy+0.55, 5.7, 0.88, size=11.5, color=TEXT_BRT)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSION & FUTURE SCOPE
# ═══════════════════════════════════════════════════════════
s9 = blank_slide()
section_tag(s9, "09  CONCLUSION")
txt(s9, "What We Achieved & What's Next",
    0.4, 0.62, 12, 0.7, size=34, bold=True, color=TEXT_BRT)
footer_bar(s9); slide_num(s9, 9)

# Left: conclusion
card(s9, 0.4, 1.5, 6.0, 5.0)
accent_bar(s9, 0.4, 1.5, 5.0, ACCENT)
txt(s9, "Conclusion", 0.65, 1.62, 5.6, 0.42, size=15, bold=True, color=ACCENT)
txts(s9, [
    "We successfully built a working ML system that:",
    "",
    "✔  Trained Linear Regression on 50 sales records",
    "✔  Achieved R² = 0.99 — near-perfect accuracy",
    "✔  Deployed via Flask with a full web dashboard",
    "✔  Added real-time USD ↔ INR currency toggle",
    "✔  Explained model decisions in plain language",
    "",
    "The project proves that even a simple algorithm",
    "like Linear Regression, when applied to clean data",
    "and deployed correctly, gives very strong results."
], 0.65, 2.12, 5.5, 4.2, size=13, color=TEXT_BRT)

# Right: future scope
card(s9, 6.8, 1.5, 6.1, 5.0)
accent_bar(s9, 6.8, 1.5, 5.0, ACCENT2)
txt(s9, "Future Scope", 7.05, 1.62, 5.7, 0.42, size=15, bold=True, color=ACCENT2)
txts(s9, [
    "This project can be extended by:",
    "",
    "→  Adding more features (season, region, day)",
    "→  Comparing with Random Forest or XGBoost",
    "→  Connecting to a live database (MySQL / Firebase)",
    "→  Adding user login and prediction history",
    "→  Hosting on the web (Render / Railway / AWS)",
    "",
    "The core architecture — ML model + REST API +",
    "interactive UI — is production-ready and can be",
    "adapted for any sales or inventory use case."
], 7.05, 2.12, 5.7, 4.2, size=13, color=TEXT_BRT)

# ── Save ────────────────────────────────────────────────────
import os
import time

# Always write to new file to avoid permission locks
timestamp = str(int(time.time()))
out = f"SalesCast_Presentation_{timestamp}.pptx"
out = os.path.abspath(out)

# Try deleting existing file first if it exists (unlock)
try:
    if os.path.exists(out):
        os.remove(out)
except:
    pass

prs.save(out)
print(f"\n✅ Successfully saved presentation at:")
print(f"   {out}")
print("\n💡 NOTE: If you had the old file open in PowerPoint, close it first before running the script.")
